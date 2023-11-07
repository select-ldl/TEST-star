import sys

from pptx import Presentation
import os
import pandas as pd

def init_fille():
    """读文件"""
    xml_files = []
    for root, dirs, files in os.walk(os.getcwd()):
        for file in files:
            if os.path.splitext(file)[1] == '.pptx':
                xml_file = os.path.join(os.getcwd(), file)  # 文件
                """处理数据"""
                xml_files.append(xml_file)
    return xml_files
def pptx_():
    data = []
    prs = Presentation(init_fille()[0])

    for slide in prs.slides:  # 遍历每页PPT
        for shape in slide.shapes:  # 遍历PPT中的每个形状
            if shape.has_text_frame:  # 判断该是否包含文本，保证有文本才提取
                for paragraph in shape.text_frame.paragraphs:  # 按文本框中的段落提取
                    data.append(paragraph.text)  # 提取一个段落的文本，就存到列表data中

    # for i in data:
    #     print(i)
    return data
def pptx11():
    data=[]
    prs = Presentation(init_fille()[0])
    a=0
    for slide in prs.slides:
        a += 1
        sheet = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    sheet.append(paragraph.text)
                    if [f"第{a}页",sheet]not in data:
                        data.append([f"第{a}页",sheet])
                    # data.append("\n")

                    # print(paragraph.text)
    return data
if __name__ == '__main__':
    # a=pptx_()

    a=pptx11()
    r = ""
    for i in a:
        a, b = i
        r = r+"\n" + str(a) + "\n"
        for m in b:
            r =r+m + "\n"
        r += "\n"
    with open("ceshi.txt","w+",encoding="utf8") as f:
        f.write(str(r))
        f.close()
    # input("\n输入任意键退出")
    # sys.exit()
    # for i in a:
    #     print(i)

